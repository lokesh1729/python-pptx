# encoding: utf-8

"""
Test suite for pptx.slide module
"""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import pytest

from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI
from pptx.package import Package
from pptx.parts.presentation import PresentationPart
from pptx.parts.slide import SlidePart
from pptx.parts.slidelayout import SlideLayoutPart
from pptx.parts.slidemaster import SlideMasterPart
from pptx.shapes.factory import SlidePlaceholders
from pptx.shapes.shapetree import SlideShapeTree
from pptx.slide import Slide, SlideLayouts, SlideMasters, Slides

from .unitutil.cxml import element, xml
from .unitutil.mock import class_mock, instance_mock, property_mock


class DescribeSlide(object):

    def it_knows_its_name(self, name_fixture):
        slide, expected_value = name_fixture
        assert slide.name == expected_value

    def it_provides_access_to_its_shapes(self, shapes_fixture):
        slide, SlideShapeTree_, shapes_ = shapes_fixture
        shapes = slide.shapes
        SlideShapeTree_.assert_called_once_with(slide)
        assert shapes is shapes_

    def it_provides_access_to_its_placeholders(self, placeholders_fixture):
        slide, SlidePlaceholders_, spTree, placeholders_ = (
            placeholders_fixture
        )
        placeholders = slide.placeholders
        SlidePlaceholders_.assert_called_once_with(spTree, slide)
        assert placeholders is placeholders_

    def it_provides_access_to_its_slide_layout(self, layout_fixture):
        slide, slide_layout_ = layout_fixture
        assert slide.slide_layout is slide_layout_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def layout_fixture(self, slide_layout_, part_prop_):
        slide = Slide(None, None)
        part_prop_.return_value.slide_layout = slide_layout_
        return slide, slide_layout_

    @pytest.fixture
    def name_fixture(self):
        sld_cxml = 'p:sld/p:cSld{name=Foobar}'
        base_slide = Slide(element(sld_cxml), None)
        return base_slide, 'Foobar'

    @pytest.fixture
    def placeholders_fixture(self, SlidePlaceholders_, placeholders_):
        sld = element('p:sld/p:cSld/p:spTree')
        slide = Slide(sld, None)
        spTree = sld.xpath('//p:spTree')[0]
        return slide, SlidePlaceholders_, spTree, placeholders_

    @pytest.fixture
    def shapes_fixture(self, SlideShapeTree_, shapes_):
        slide = Slide(None, None)
        return slide, SlideShapeTree_, shapes_

    # fixture components -----------------------------------

    @pytest.fixture
    def part_prop_(self, request):
        return property_mock(request, Slide, 'part')

    @pytest.fixture
    def placeholders_(self, request):
        return instance_mock(request, SlidePlaceholders)

    @pytest.fixture
    def SlidePlaceholders_(self, request, placeholders_):
        return class_mock(
            request, 'pptx.slide.SlidePlaceholders',
            return_value=placeholders_
        )

    @pytest.fixture
    def SlideShapeTree_(self, request, shapes_):
        return class_mock(
            request, 'pptx.slide.SlideShapeTree', return_value=shapes_
        )

    @pytest.fixture
    def shapes_(self, request):
        return instance_mock(request, SlideShapeTree)

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayoutPart)


class DescribeSlides(object):

    def it_supports_indexed_access(self, getitem_fixture):
        slides, prs_part_, slide_, rId = getitem_fixture
        slide = slides[0]
        prs_part_.related_parts.__getitem__.assert_called_once_with(rId)
        assert slide is slide_

    def it_raises_on_slide_index_out_of_range(self, getitem_raises_fixture):
        slides = getitem_raises_fixture
        with pytest.raises(IndexError):
            slides[2]

    def it_knows_the_index_of_a_slide_it_contains(self, index_fixture):
        slides, slide, expected_value = index_fixture
        index = slides.index(slide)
        assert index == expected_value

    def it_raises_on_slide_not_in_collection(self, raises_fixture):
        slides, slide = raises_fixture
        with pytest.raises(ValueError):
            slides.index(slide)

    def it_can_iterate_its_slides(self, iter_fixture):
        slides, expected_value = iter_fixture
        slide_lst = [s for s in slides]
        assert slide_lst == expected_value

    def it_supports_len(self, len_fixture):
        slides, expected_value = len_fixture
        assert len(slides) == expected_value

    def it_can_add_a_new_slide(self, add_fixture):
        slides, slide_layout_, SlidePart_, part_name = add_fixture[:4]
        package_, slide_part_, slide_, expected_xml = add_fixture[4:]

        slide = slides.add_slide(slide_layout_)

        SlidePart_.new.assert_called_once_with(
            slide_layout_, part_name, package_
        )
        slides.part.relate_to.assert_called_once_with(slide_part_, RT.SLIDE)
        assert slides._sldIdLst.xml == expected_xml
        assert slide is slide_

    def it_assigns_partnames_to_its_slides_to_help(self, rename_fixture):
        slides, slide_lst, expected_names = rename_fixture
        slides.rename_slides()
        assert [s.partname for s in slide_lst] == expected_names

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def add_fixture(
            self, part_prop_, slide_layout_, SlidePart_, package_,
            slide_part_):
        slides = Slides(element('p:sldIdLst/p:sldId{r:id=rId1}'), None)
        part_name = PackURI('/ppt/slides/slide2.xml')
        slide_ = slide_part_.slide
        expected_xml = xml(
            'p:sldIdLst/(p:sldId{r:id=rId1},p:sldId{r:id=rId2,id=256})'
        )
        SlidePart_.new.return_value = slide_part_
        part_prop_.return_value.package = package_
        part_prop_.return_value.relate_to.return_value = 'rId2'
        return (
            slides, slide_layout_, SlidePart_, part_name, package_,
            slide_part_, slide_, expected_xml
        )

    @pytest.fixture
    def getitem_fixture(self, prs_part_, slide_part_, part_prop_):
        sldIdLst = element('p:sldIdLst/p:sldId{r:id=rId1}')
        slides = Slides(sldIdLst, None)
        part_prop_.return_value = prs_part_
        prs_part_.related_parts.__getitem__.return_value = slide_part_
        slide_ = slide_part_.slide
        return slides, prs_part_, slide_, 'rId1'

    @pytest.fixture
    def getitem_raises_fixture(self):
        sldIdLst = element('p:sldIdLst/p:sldId{r:id=rId1}')
        slides = Slides(sldIdLst, None)
        return slides

    @pytest.fixture(params=[0, 1])
    def index_fixture(self, request, part_prop_):
        idx = request.param
        sldIdLst = element('p:sldIdLst/(p:sldId{r:id=a},p:sldId{r:id=b})')
        slides = Slides(sldIdLst, None)
        _slides = [
            SlidePart(None, None, element('p:sld')),
            SlidePart(None, None, element('p:sld'))
        ]
        part_prop_.return_value.related_parts.__getitem__.side_effect = (
            _slides
        )
        return slides, _slides[idx], idx

    @pytest.fixture
    def iter_fixture(self, part_prop_):
        sldIdLst = element('p:sldIdLst/(p:sldId{r:id=a},p:sldId{r:id=b})')
        slides = Slides(sldIdLst, None)
        related_parts_ = part_prop_.return_value.related_parts
        slide_part = related_parts_.__getitem__.return_value
        _slides = [slide_part.slide, slide_part.slide]
        return slides, _slides

    @pytest.fixture(params=[
        ('p:sldIdLst',                                   0),
        ('p:sldIdLst/p:sldId{r:id=a}',                   1),
        ('p:sldIdLst/(p:sldId{r:id=a},p:sldId{r:id=b})', 2),
    ])
    def len_fixture(self, request):
        sldIdLst_cxml, expected_value = request.param
        slides = Slides(element(sldIdLst_cxml), None)
        return slides, expected_value

    @pytest.fixture
    def raises_fixture(self):
        slides = Slides(element('p:sldIdLst'), None)
        slide = SlidePart(None, None, element('p:sld'))
        return slides, slide

    @pytest.fixture
    def rename_fixture(self, part_prop_):
        sldIdLst = element('p:sldIdLst/(p:sldId{r:id=a},p:sldId{r:id=b})')
        slides = Slides(sldIdLst, None)
        related_parts_ = part_prop_.return_value.related_parts
        slide_part = related_parts_.__getitem__.return_value
        _slides = [slide_part.slide, slide_part.slide]
        expected_names = [
            PackURI('/ppt/slides/slide2.xml'),
            PackURI('/ppt/slides/slide2.xml'),
        ]
        return slides, _slides, expected_names

    # fixture components ---------------------------------------------

    @pytest.fixture
    def package_(self, request):
        return instance_mock(request, Package)

    @pytest.fixture
    def part_prop_(self, request):
        return property_mock(request, Slides, 'part')

    @pytest.fixture
    def prs_part_(self, request):
        return instance_mock(request, PresentationPart)

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayoutPart)

    @pytest.fixture
    def SlidePart_(self, request, slide_part_):
        return class_mock(
            request, 'pptx.slide.SlidePart', return_value=slide_part_
        )

    @pytest.fixture
    def slide_part_(self, request):
        return instance_mock(request, SlidePart)


class DescribeSlideLayouts(object):

    def it_supports_len(self, len_fixture):
        slide_layouts, expected_value = len_fixture
        assert len(slide_layouts) == expected_value

    def it_can_iterate_its_slide_layouts(self, iter_fixture):
        slide_layouts, expected_value = iter_fixture
        assert [s for s in slide_layouts] == expected_value

    def it_supports_indexed_access(self, getitem_fixture):
        slide_layouts, part_, slide_layout_, rId = getitem_fixture
        slide_layout = slide_layouts[0]
        part_.related_parts.__getitem__.assert_called_once_with(rId)
        assert slide_layout is slide_layout_

    def it_raises_on_index_out_of_range(self, getitem_raises_fixture):
        slides = getitem_raises_fixture
        with pytest.raises(IndexError):
            slides[1]

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def getitem_fixture(self, part_, slide_layout_, part_prop_):
        slide_layouts = SlideLayouts(
            element('p:sldLayoutIdLst/p:sldLayoutId{r:id=rId1}'), None
        )
        part_prop_.return_value = part_
        part_.related_parts.__getitem__.return_value = slide_layout_
        return slide_layouts, part_, slide_layout_, 'rId1'

    @pytest.fixture
    def getitem_raises_fixture(self, part_prop_):
        return SlideLayouts(
            element('p:sldLayoutIdLst/p:sldLayoutId{r:id=rId1}'), None
        )

    @pytest.fixture
    def iter_fixture(self, part_prop_):
        sldLayoutIdLst = element(
            'p:sldLayoutIdLst/(p:sldLayoutId{r:id=a},p:sldLayoutId{r:id=b})'
        )
        slide_layouts = SlideLayouts(sldLayoutIdLst, None)
        _slide_layouts = [
            SlideLayoutPart(None, None, element('p:sldLayout')),
            SlideLayoutPart(None, None, element('p:sldLayout')),
        ]
        part_prop_.return_value.related_parts.__getitem__.side_effect = (
            _slide_layouts
        )
        return slide_layouts, _slide_layouts

    @pytest.fixture(params=[
        ('p:sldLayoutIdLst',                               0),
        ('p:sldLayoutIdLst/p:sldLayoutId',                 1),
        ('p:sldLayoutIdLst/(p:sldLayoutId,p:sldLayoutId)', 2),
    ])
    def len_fixture(self, request):
        sldLayoutIdLst_cxml, expected_value = request.param
        slide_layouts = SlideLayouts(element(sldLayoutIdLst_cxml), None)
        return slide_layouts, expected_value

    # fixture components ---------------------------------------------

    @pytest.fixture
    def part_(self, request):
        return instance_mock(request, SlideMasterPart)

    @pytest.fixture
    def part_prop_(self, request):
        return property_mock(request, SlideLayouts, 'part')

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayoutPart)


class DescribeSlideMasters(object):

    def it_knows_how_many_masters_it_contains(self, len_fixture):
        slide_masters, expected_value = len_fixture
        assert len(slide_masters) == expected_value

    def it_can_iterate_the_slide_masters(self, iter_fixture):
        slide_masters, expected_values = iter_fixture
        assert [sm for sm in slide_masters] == expected_values

    def it_supports_indexed_access(self, getitem_fixture):
        slide_masters, part_, slide_master_, rId = getitem_fixture
        slide_master = slide_masters[0]
        part_.related_parts.__getitem__.assert_called_once_with(rId)
        assert slide_master is slide_master_

    def it_raises_on_index_out_of_range(self, getitem_raises_fixture):
        slides = getitem_raises_fixture
        with pytest.raises(IndexError):
            slides[1]

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def getitem_fixture(self, part_, slide_master_, part_prop_):
        slide_masters = SlideMasters(
            element('p:sldMasterIdLst/p:sldMasterId{r:id=rId1}'), None
        )
        part_prop_.return_value = part_
        part_.related_parts.__getitem__.return_value = slide_master_
        return slide_masters, part_, slide_master_, 'rId1'

    @pytest.fixture
    def getitem_raises_fixture(self, part_prop_):
        return SlideMasters(
            element('p:sldMasterIdLst/p:sldMasterId{r:id=rId1}'), None
        )

    @pytest.fixture
    def iter_fixture(self, part_prop_):
        sldMasterIdLst = element(
            'p:sldMasterIdLst/(p:sldMasterId{r:id=a},p:sldMasterId{r:id=b})'
        )
        slide_masters = SlideMasters(sldMasterIdLst, None)
        _slide_masters = [
            SlideMasterPart(None, None, element('p:sldMaster')),
            SlideMasterPart(None, None, element('p:sldMaster'))
        ]
        part_prop_.return_value.related_parts.__getitem__.side_effect = (
            _slide_masters
        )
        return slide_masters, _slide_masters

    @pytest.fixture(params=[
        ('p:sldMasterIdLst',                               0),
        ('p:sldMasterIdLst/p:sldMasterId',                 1),
        ('p:sldMasterIdLst/(p:sldMasterId,p:sldMasterId)', 2),
    ])
    def len_fixture(self, request):
        sldMasterIdLst_cxml, expected_value = request.param
        slide_masters = SlideMasters(element(sldMasterIdLst_cxml), None)
        return slide_masters, expected_value

    # fixture components ---------------------------------------------

    @pytest.fixture
    def part_(self, request):
        return instance_mock(request, PresentationPart)

    @pytest.fixture
    def part_prop_(self, request):
        return property_mock(request, SlideMasters, 'part')

    @pytest.fixture
    def slide_master_(self, request):
        return instance_mock(request, SlideMasterPart)
